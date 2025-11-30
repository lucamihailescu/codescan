import os
import hashlib
import mimetypes
import pickle
from concurrent.futures import ThreadPoolExecutor, as_completed
from threading import Lock
from sqlalchemy.orm import Session
from sklearn.feature_extraction.text import TfidfVectorizer, HashingVectorizer
from models import IndexedFile, IndexOperation
from datetime import datetime, timezone
from progress_store import progress_store
from similarity_config import similarity_config_store
from storage_config import storage_config_store
from storage_factory import get_storage_backend
from storage_interface import StorageBackendInterface
from ignored_files_config import ignored_files_store
from typing import Optional, Any

# Optional document extraction libraries
# These are typed as Any to satisfy the type checker when conditionally imported
DocxDocument: Any = None
pypdf: Any = None
openpyxl: Any = None
Presentation: Any = None

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pypdf
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

def get_vectorizer():
    """Get a TF-IDF vectorizer configured based on current similarity settings"""
    config = similarity_config_store.config
    return TfidfVectorizer(
        max_features=config.n_features,
        ngram_range=(config.ngram_range_min, config.ngram_range_max),
        use_idf=config.use_idf,
        sublinear_tf=config.sublinear_tf,
        max_df=config.max_df,
        min_df=config.min_df,
        lowercase=True,
        strip_accents='unicode',
        stop_words='english',  # Remove common English stop words
    )

def count_files(directory: str) -> int:
    """Count total files in directory for progress tracking (excludes ignored files)"""
    count = 0
    
    def walk_error_handler(error: OSError):
        """Handle errors during directory traversal (e.g., permission denied)"""
        print(f"Warning: Cannot access {error.filename}: {error.strerror}")
    
    for root, dirs, files in os.walk(directory, onerror=walk_error_handler):
        for filename in files:
            if not ignored_files_store.should_ignore(filename):
                count += 1
    return count

def get_file_hash(filepath: str) -> str:
    """Compute SHA256 hash of file. Raises PermissionError if access denied."""
    sha256_hash = hashlib.sha256()
    try:
        with open(filepath, "rb") as f:
            # Read and update hash string value in blocks of 4K
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)
        return sha256_hash.hexdigest()
    except PermissionError:
        raise PermissionError(f"Access denied: {filepath}")
    except OSError as e:
        raise OSError(f"Cannot read file {filepath}: {e.strerror}")


# ============ Document Type Detection ============

# Supported document extensions (require special extraction)
DOCUMENT_EXTENSIONS = {
    '.docx': 'word',
    '.doc': 'word_legacy',
    '.pdf': 'pdf',
    '.xlsx': 'excel',
    '.xls': 'excel_legacy',
    '.pptx': 'powerpoint',
    '.ppt': 'powerpoint_legacy',
    '.odt': 'odt',
    '.rtf': 'rtf',
}

# Source code and config file extensions (always treat as text)
CODE_EXTENSIONS = {
    # Python
    '.py', '.pyw', '.pyi', '.pyx',
    # Java
    '.java', '.jar', '.class',
    # C/C++
    '.c', '.h', '.cpp', '.hpp', '.cc', '.hh', '.cxx', '.hxx',
    # C#
    '.cs', '.csx',
    # JavaScript/TypeScript
    '.js', '.jsx', '.ts', '.tsx', '.mjs', '.cjs',
    # Web
    '.html', '.htm', '.css', '.scss', '.sass', '.less',
    # Data/Config
    '.json', '.xml', '.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf',
    # Shell/Scripts
    '.sh', '.bash', '.zsh', '.fish', '.ps1', '.bat', '.cmd',
    # Ruby
    '.rb', '.rake', '.gemspec',
    # Go
    '.go',
    # Rust
    '.rs',
    # Swift/Objective-C
    '.swift', '.m', '.mm',
    # Kotlin
    '.kt', '.kts',
    # Scala
    '.scala', '.sc',
    # PHP
    '.php', '.phtml',
    # SQL
    '.sql',
    # Markdown/Text
    '.md', '.markdown', '.txt', '.text', '.rst', '.asciidoc',
    # Other
    '.r', '.R', '.jl', '.lua', '.pl', '.pm', '.groovy', '.gradle',
    '.dockerfile', '.makefile', '.cmake',
}

def get_file_type(filepath: str) -> str:
    """
    Determine file type for text extraction.
    Returns: 'text', 'word', 'pdf', 'excel', 'powerpoint', or 'binary'
    """
    ext = os.path.splitext(filepath)[1].lower()
    filename = os.path.basename(filepath).lower()
    
    # Check for known document formats (require special extraction)
    if ext in DOCUMENT_EXTENSIONS:
        return DOCUMENT_EXTENSIONS[ext]
    
    # Check for known code/text file extensions
    if ext in CODE_EXTENSIONS:
        return 'text'
    
    # Check for common config files without extensions
    if filename in ('dockerfile', 'makefile', 'gemfile', 'rakefile', 'procfile',
                    '.gitignore', '.dockerignore', '.env', '.editorconfig'):
        return 'text'
    
    # Check MIME type for text files
    mime_type, _ = mimetypes.guess_type(filepath)
    if mime_type and mime_type.startswith('text'):
        return 'text'
    
    # Try reading as UTF-8 text
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            f.read(1024)
        return 'text'
    except (UnicodeDecodeError, PermissionError):
        pass
    
    return 'binary'


def is_text_file(filepath: str) -> bool:
    """Check if file contains extractable text content."""
    file_type = get_file_type(filepath)
    return file_type != 'binary'


# ============ Document Text Extraction ============

def extract_text_from_docx(filepath: str) -> str:
    """Extract text from Word .docx files."""
    if not DOCX_AVAILABLE:
        print(f"python-docx not installed, cannot read {filepath}")
        return ""
    try:
        doc = DocxDocument(filepath)
        paragraphs = [para.text for para in doc.paragraphs]
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    paragraphs.append(cell.text)
        return "\n".join(paragraphs)
    except Exception as e:
        print(f"Error extracting text from {filepath}: {e}")
        return ""


def extract_text_from_pdf(filepath: str) -> str:
    """Extract text from PDF files."""
    if not PDF_AVAILABLE:
        print(f"pypdf not installed, cannot read {filepath}")
        return ""
    try:
        text_parts = []
        with open(filepath, 'rb') as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting text from {filepath}: {e}")
        return ""


def extract_text_from_xlsx(filepath: str) -> str:
    """Extract text from Excel .xlsx files."""
    if not XLSX_AVAILABLE:
        print(f"openpyxl not installed, cannot read {filepath}")
        return ""
    try:
        workbook = openpyxl.load_workbook(filepath, data_only=True)
        text_parts = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        text_parts.append(str(cell.value))
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting text from {filepath}: {e}")
        return ""


def extract_text_from_pptx(filepath: str) -> str:
    """Extract text from PowerPoint .pptx files."""
    if not PPTX_AVAILABLE:
        print(f"python-pptx not installed, cannot read {filepath}")
        return ""
    try:
        prs = Presentation(filepath)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(getattr(shape, "text", ""))
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting text from {filepath}: {e}")
        return ""


def extract_text_from_file(filepath: str) -> str:
    """
    Extract text content from a file based on its type.
    Returns extracted text or empty string if extraction fails.
    """
    file_type = get_file_type(filepath)
    
    if file_type == 'text':
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            print(f"Error reading text file {filepath}: {e}")
            return ""
    
    elif file_type == 'word':
        return extract_text_from_docx(filepath)
    
    elif file_type == 'pdf':
        return extract_text_from_pdf(filepath)
    
    elif file_type == 'excel':
        return extract_text_from_xlsx(filepath)
    
    elif file_type == 'powerpoint':
        return extract_text_from_pptx(filepath)
    
    elif file_type in ('word_legacy', 'excel_legacy', 'powerpoint_legacy', 'odt', 'rtf'):
        # Legacy formats not yet supported
        print(f"Legacy format {file_type} not supported for {filepath}")
        return ""
    
    return ""

from typing import Optional

def compute_vector(filepath: Optional[str] = None, content: Optional[str] = None):
    """
    Compute TF-IDF vector for file content.
    Uses HashingVectorizer for consistent hashing across different runs.
    Either filepath or content must be provided.
    """
    try:
        if content is None:
            if filepath is None:
                return None
            # Use document extraction to get text content
            content = extract_text_from_file(filepath)
            if not content:
                return None
        
        config = similarity_config_store.config
        
        # Skip very short content
        if len(content.strip()) < config.min_content_length:
            return None
        
        # Use HashingVectorizer for consistent hashing (stateless, no fitting required)
        # This ensures vectors computed at different times are comparable
        local_vectorizer = HashingVectorizer(
            n_features=config.n_features,
            ngram_range=(config.ngram_range_min, config.ngram_range_max),
            alternate_sign=False,
            norm='l2',
            lowercase=True,
            strip_accents='unicode',
            stop_words='english',
        )
        vector = local_vectorizer.transform([content])
        return pickle.dumps(vector)
    except Exception as e:
        print(f"Error computing vector for {filepath}: {e}")
        return None


def compute_vector_from_content(content: str):
    """Compute vector directly from content string"""
    return compute_vector(content=content)


def _index_single_file_with_storage(filepath: str, storage: StorageBackendInterface) -> bool:
    """Index a single file using storage abstraction. Returns True if file was indexed/updated.
    
    Raises:
        PermissionError: If the file cannot be accessed due to permissions.
        OSError: If the file cannot be read for other OS-related reasons.
    """
    try:
        stat = os.stat(filepath)
    except PermissionError:
        raise PermissionError(f"Access denied: {filepath}")
    except OSError as e:
        raise OSError(f"Cannot access file {filepath}: {e.strerror}")
    
    last_modified = stat.st_mtime
    
    # Check if already indexed and not modified
    existing = storage.get_indexed_file_by_path(filepath)
    if existing is not None and float(existing.last_modified) == last_modified:
        # Also check if vector is missing for text files - if so, re-index
        if is_text_file(filepath) and existing.vector is None:
            print(f"Re-indexing {filepath} - missing vector")
        else:
            return False  # Skip, not modified and has vector

    # Compute Hash (may raise PermissionError/OSError)
    file_hash = get_file_hash(filepath)
    
    # Compute Vector if text
    vector_blob = None
    if is_text_file(filepath):
        vector_blob = compute_vector(filepath)

    # Add or update via storage abstraction
    storage.add_or_update_indexed_file(
        path=filepath,
        filename=os.path.basename(filepath),
        file_hash=file_hash,
        vector=vector_blob,
        last_modified=last_modified
    )
    
    return True


def index_directory(directory: str, db: Session):
    """Original index_directory for backward compatibility (no progress tracking)"""
    # Use storage abstraction if Redis is configured, otherwise use direct SQLAlchemy
    if storage_config_store.is_redis():
        storage = get_storage_backend()
        try:
            for root, dirs, files in os.walk(directory):
                for file in files:
                    if ignored_files_store.should_ignore(file):
                        continue
                    filepath = os.path.join(root, file)
                    try:
                        _index_single_file_with_storage(filepath, storage)
                        print(f"Indexed: {filepath}")
                    except Exception as e:
                        print(f"Failed to index {filepath}: {e}")
        finally:
            storage.close()
    else:
        # Original SQLAlchemy implementation
        for root, dirs, files in os.walk(directory):
            for file in files:
                if ignored_files_store.should_ignore(file):
                    continue
                filepath = os.path.join(root, file)
                try:
                    _index_single_file(filepath, db)
                    print(f"Indexed: {filepath}")
                except Exception as e:
                    print(f"Failed to index {filepath}: {e}")
                    db.rollback()


def _index_single_file(filepath: str, db: Session) -> bool:
    """Index a single file. Returns True if file was indexed/updated.
    
    Raises:
        PermissionError: If the file cannot be accessed due to permissions.
        OSError: If the file cannot be read for other OS-related reasons.
    """
    try:
        stat = os.stat(filepath)
    except PermissionError:
        raise PermissionError(f"Access denied: {filepath}")
    except OSError as e:
        raise OSError(f"Cannot access file {filepath}: {e.strerror}")
    
    last_modified = stat.st_mtime
    
    # Check if already indexed and not modified
    existing = db.query(IndexedFile).filter(IndexedFile.path == filepath).first()
    if existing is not None and float(existing.last_modified) == last_modified:
        # Force re-index if vector is missing but file should have one (text file)
        if existing.vector is None and is_text_file(filepath):
            pass  # Continue to re-index
        else:
            return False  # Skip, not modified

    # Compute Hash (may raise PermissionError/OSError)
    file_hash = get_file_hash(filepath)
    
    # Compute Vector if text
    vector_blob = None
    if is_text_file(filepath):
        vector_blob = compute_vector(filepath)

    if existing:
        existing.file_hash = file_hash
        existing.vector = vector_blob
        existing.last_modified = last_modified
        existing.indexed_at = datetime.now(timezone.utc)
    else:
        new_file = IndexedFile(
            path=filepath,
            filename=os.path.basename(filepath),
            file_hash=file_hash,
            vector=vector_blob,
            last_modified=last_modified
        )
        db.add(new_file)
    
    db.commit()
    return True


def _collect_files(directory: str) -> list:
    """Collect all file paths from directory (excludes ignored files)"""
    files = []
    skipped_dirs = []
    
    def walk_error_handler(error: OSError):
        """Handle errors during directory traversal (e.g., permission denied)"""
        skipped_dirs.append(error.filename)
        print(f"Warning: Cannot access directory {error.filename}: {error.strerror}")
    
    for root, dirs, filenames in os.walk(directory, onerror=walk_error_handler):
        for filename in filenames:
            if not ignored_files_store.should_ignore(filename):
                files.append(os.path.join(root, filename))
    
    if skipped_dirs:
        print(f"Skipped {len(skipped_dirs)} inaccessible directories")
    return files


def _index_file_worker(filepath: str, use_redis: bool, db: Optional[Session] = None) -> tuple:
    """
    Worker function for parallel indexing.
    Returns (filepath, was_indexed, error_msg)
    """
    try:
        if use_redis:
            # Each thread gets its own storage connection
            storage = get_storage_backend()
            try:
                was_indexed = _index_single_file_with_storage(filepath, storage)
                return (filepath, was_indexed, None)
            finally:
                storage.close()
        else:
            # For SQLite, we need to be careful with session management
            # Create a new session for this thread
            from database import SessionLocal
            thread_db = SessionLocal()
            try:
                was_indexed = _index_single_file(filepath, thread_db)
                thread_db.commit()  # Commit the changes
                return (filepath, was_indexed, None)
            except Exception as e:
                thread_db.rollback()
                raise
            finally:
                thread_db.close()
    except Exception as e:
        return (filepath, False, str(e))


def index_directory_with_id(directory: str, db: Session, index_id: str):
    """Index directory with progress tracking via index_id"""
    print(f"Starting indexing {index_id} on {directory}")
    
    # Create index operation record
    index_op = IndexOperation(
        index_id=index_id,
        directory_path=directory,
        status="running",
        started_at=datetime.now(timezone.utc)
    )
    db.add(index_op)
    db.commit()
    
    # Check threading configuration
    threading_config = storage_config_store.config.threading_config
    use_threading = threading_config.enabled
    
    if use_threading:
        return _index_directory_parallel(directory, db, index_id, threading_config)
    else:
        return _index_directory_sequential(directory, db, index_id)


def _index_directory_sequential(directory: str, db: Session, index_id: str):
    """Original sequential indexing implementation"""
    # Update progress to counting status
    progress_store.update_task(index_id, status="counting")
    
    # Count total files first
    total_files = count_files(directory)
    progress_store.update_task(index_id, status="processing", total_files=total_files)
    
    files_processed = 0
    files_indexed = 0
    files_access_denied = 0
    
    # Determine which indexing method to use based on storage backend
    use_redis = storage_config_store.is_redis()
    storage = None
    
    if use_redis:
        storage = get_storage_backend()
    
    def walk_error_handler(error: OSError):
        """Handle errors during directory traversal (e.g., permission denied)"""
        print(f"Warning: Cannot access directory {error.filename}: {error.strerror}")
    
    cancelled = False
    try:
        for root, dirs, files in os.walk(directory, onerror=walk_error_handler):
            for file in files:
                # Skip ignored files
                if ignored_files_store.should_ignore(file):
                    continue
                    
                # Check if cancelled
                if progress_store.is_cancelled(index_id):
                    print(f"Indexing cancelled: {index_id}")
                    cancelled = True
                    break
                    
                filepath = os.path.join(root, file)
                files_processed += 1
                
                # Update progress
                progress_store.update_task(
                    index_id,
                    files_processed=files_processed,
                    current_file=filepath
                )
                
                try:
                    if use_redis and storage:
                        was_indexed = _index_single_file_with_storage(filepath, storage)
                    else:
                        was_indexed = _index_single_file(filepath, db)
                    
                    if was_indexed:
                        files_indexed += 1
                        progress_store.update_task(index_id, files_indexed=files_indexed)
                        print(f"Indexed: {filepath}")
                    else:
                        print(f"Skipped (unchanged): {filepath}")
                except PermissionError as e:
                    files_access_denied += 1
                    print(f"Access denied: {filepath}")
                except OSError as e:
                    files_access_denied += 1
                    print(f"Cannot access file {filepath}: {e}")
                except Exception as e:
                    print(f"Failed to index {filepath}: {e}")
                    if not use_redis:
                        db.rollback()
            
            # Also check after inner loop
            if cancelled:
                break
    finally:
        if storage:
            storage.close()
    
    # Determine final status
    final_status = "cancelled" if cancelled else "completed"
    
    # Mark indexing as completed/cancelled
    progress_store.update_task(
        index_id,
        status=final_status,
        completed_at=datetime.now(),
        current_file=""
    )
    progress_store.clear_cancelled(index_id)
    
    # Update index operation record
    from database import SessionLocal
    update_db = SessionLocal()
    try:
        index_op = update_db.query(IndexOperation).filter(IndexOperation.index_id == index_id).first()
        if index_op:
            index_op.status = final_status
            index_op.total_files = total_files
            index_op.files_indexed = files_indexed
            index_op.files_skipped = files_processed - files_indexed
            index_op.completed_at = datetime.now(timezone.utc)
            update_db.commit()
    finally:
        update_db.close()
    
    if cancelled:
        print(f"Indexing cancelled: {files_indexed} files indexed out of {files_processed} processed ({total_files} total)")
    else:
        access_msg = f", {files_access_denied} access denied" if files_access_denied > 0 else ""
        print(f"Indexing completed: {files_indexed} files indexed out of {total_files} total{access_msg}")
    return index_id


def _index_directory_parallel(directory: str, db: Session, index_id: str, threading_config):
    """Parallel indexing implementation using ThreadPoolExecutor"""
    from storage_config import ThreadingConfig
    
    # Update progress to counting status
    progress_store.update_task(index_id, status="counting")
    
    # Collect all files first
    all_files = _collect_files(directory)
    total_files = len(all_files)
    progress_store.update_task(index_id, status="processing", total_files=total_files)
    
    files_processed = 0
    files_indexed = 0
    files_access_denied = 0
    progress_lock = Lock()
    
    use_redis = storage_config_store.is_redis()
    max_workers = threading_config.max_workers
    
    print(f"Starting parallel indexing with {max_workers} workers")
    
    def update_progress(filepath: str, was_indexed: bool, error: Optional[str] = None):
        nonlocal files_processed, files_indexed, files_access_denied
        with progress_lock:
            files_processed += 1
            if was_indexed:
                files_indexed += 1
                print(f"Indexed: {filepath}")
            elif error:
                # Check if it's an access denied error
                if "Access denied" in error or "Permission denied" in error:
                    files_access_denied += 1
                    print(f"Access denied: {filepath}")
                else:
                    print(f"Failed to index {filepath}: {error}")
            else:
                print(f"Skipped (unchanged): {filepath}")
            
            progress_store.update_task(
                index_id,
                files_processed=files_processed,
                files_indexed=files_indexed,
                current_file=filepath
            )
    
    cancelled = False
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        # Submit all indexing tasks
        futures = {
            executor.submit(_index_file_worker, filepath, use_redis, db): filepath 
            for filepath in all_files
        }
        
        # Process results as they complete
        for future in as_completed(futures):
            # Check if cancelled
            if progress_store.is_cancelled(index_id):
                print(f"Indexing cancelled: {index_id}")
                cancelled = True
                # Cancel remaining futures
                for f in futures:
                    f.cancel()
                break
            
            filepath, was_indexed, error = future.result()
            update_progress(filepath, was_indexed, error)
    
    # Determine final status
    final_status = "cancelled" if cancelled else "completed"
    
    # Mark indexing as completed/cancelled
    progress_store.update_task(
        index_id,
        status=final_status,
        completed_at=datetime.now(),
        current_file=""
    )
    progress_store.clear_cancelled(index_id)
    
    # Update index operation record
    from database import SessionLocal
    update_db = SessionLocal()
    try:
        index_op = update_db.query(IndexOperation).filter(IndexOperation.index_id == index_id).first()
        if index_op:
            index_op.status = final_status
            index_op.total_files = total_files
            index_op.files_indexed = files_indexed
            index_op.files_skipped = files_processed - files_indexed
            index_op.completed_at = datetime.now(timezone.utc)
            update_db.commit()
    finally:
        update_db.close()
    
    if cancelled:
        print(f"Parallel indexing cancelled: {files_indexed} files indexed out of {files_processed} processed ({total_files} total)")
    else:
        access_msg = f", {files_access_denied} access denied" if files_access_denied > 0 else ""
        print(f"Parallel indexing completed: {files_indexed} files indexed out of {total_files} total{access_msg}")
    return index_id
